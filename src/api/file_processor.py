"""
File Processor — Watchdog-Based Workspace File Watcher
=======================================================

Monitors the workspace directory for new or modified files and automatically
extracts their text content into ``.processed/`` for LLM consumption.

Supported formats:
- Documents: PDF, DOCX (via Docling — layout-aware, table extraction)
- Presentations: PPTX, EPUB, RTF
- Data: CSV, XLSX, JSON, XML, YAML, TOML, INI
- Web: HTML, Markdown
- Database: SQLite
- Archives: ZIP, TAR, GZ
- Source code: Python, JS, TS, Java, C++, Go, Rust, Ruby, PHP
- Logs and plain text (fallback)

The watcher runs in a background thread and notifies the frontend via a
callback (typically a WebSocket broadcast) when processing completes.
"""

import os
import time
import threading
import logging
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

# Configure simple logging
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s"
)
logger = logging.getLogger(__name__)

# Files explicitly uploaded via frontend; bypasses Eco-Mode restrictions
FORCE_PROCESS_FILES = set()


class FileWatcherHandler(FileSystemEventHandler):
    """Handles filesystem events, queuing files for text extraction.

    On file create/modify, waits briefly for write locks to release,
    then dispatches to a format-specific processor. Output goes to
    ``{workspace}/.processed/{filename}.txt`` (or ``.md`` for tables).

    PDF and DOCX extraction uses Docling when available (layout-aware
    with table structure detection), falling back to PyMuPDF and
    python-docx respectively.

    Args:
        workspace_dir: Absolute path to the watched workspace directory.
        on_processed_callback: Optional ``(filename, status)`` callback
            invoked after each file is processed (or on error).
    """

    def __init__(self, workspace_dir, on_processed_callback=None):
        self.workspace_dir = os.path.abspath(workspace_dir)
        self.processed_dir = os.path.join(self.workspace_dir, ".processed")
        os.makedirs(self.processed_dir, exist_ok=True)
        self._docling_converter = None
        self.on_processed_callback = on_processed_callback  # callback(filename)

    def on_created(self, event):
        if event.is_directory:
            return
        logger.info(f"[Watcher] File created: {event.src_path}")
        self._trigger_processing(event.src_path)

    def on_modified(self, event):
        if event.is_directory:
            return
        logger.info(f"[Watcher] File modified: {event.src_path}")
        self._trigger_processing(event.src_path)

    def on_deleted(self, event):
        if event.is_directory:
            return
        logger.info(f"[Watcher] File deleted: {event.src_path}")
        # Skip hidden files
        if os.path.basename(event.src_path).startswith("."):
            return
        # Skip if within the .processed directory
        if event.src_path.startswith(self.processed_dir):
            return
        from src.memory.vector_lifecycle import VectorLifecycleManager

        filename = os.path.basename(event.src_path)
        import asyncio

        try:
            loop = asyncio.get_event_loop()
            if loop.is_running():
                loop.create_task(
                    VectorLifecycleManager.on_file_deleted(self.workspace_dir, filename)
                )
            else:
                loop.run_until_complete(
                    VectorLifecycleManager.on_file_deleted(self.workspace_dir, filename)
                )
        except RuntimeError:
            asyncio.run(
                VectorLifecycleManager.on_file_deleted(self.workspace_dir, filename)
            )

    def on_moved(self, event):
        if event.is_directory:
            return
        logger.info(f"[Watcher] File moved: {event.src_path} -> {event.dest_path}")
        # Skip hidden files
        if os.path.basename(event.src_path).startswith("."):
            return
        if event.src_path.startswith(self.processed_dir):
            return
        from src.memory.vector_lifecycle import VectorLifecycleManager

        old_filename = os.path.basename(event.src_path)
        new_filename = os.path.basename(event.dest_path)
        import asyncio

        try:
            loop = asyncio.get_event_loop()
            if loop.is_running():
                loop.create_task(
                    VectorLifecycleManager.on_file_renamed(
                        self.workspace_dir, old_filename, new_filename
                    )
                )
            else:
                loop.run_until_complete(
                    VectorLifecycleManager.on_file_renamed(
                        self.workspace_dir, old_filename, new_filename
                    )
                )
        except RuntimeError:
            asyncio.run(
                VectorLifecycleManager.on_file_renamed(
                    self.workspace_dir, old_filename, new_filename
                )
            )

    def _trigger_processing(self, filepath):
        # Skip hidden files and already processed files to prevent loops
        if os.path.basename(filepath).startswith("."):
            return
        # Skip if within the .processed directory
        if filepath.startswith(self.processed_dir):
            return

        try:
            from src.api.power_monitor import ECO_MODE

            if ECO_MODE and filepath not in FORCE_PROCESS_FILES:
                logger.info(
                    f"[Watcher] Eco-Mode active: Skipping background extraction for {filepath}"
                )
                return
            if filepath in FORCE_PROCESS_FILES:
                FORCE_PROCESS_FILES.discard(filepath)
        except ImportError:
            pass

        # Run processing in a background thread with slight delay for file locked operations (uploads)
        threading.Thread(
            target=self._delayed_process, args=(filepath,), daemon=True
        ).start()

    def _delayed_process(self, filepath, delay=1.0):
        """Wait for potential write operations to complete before processing."""
        time.sleep(delay)
        try:
            # Check file still exists after delay
            if not os.path.exists(filepath):
                return

            self.process_file(filepath)
        except Exception as e:
            logger.error(f"[Watcher] Error in delayed process for {filepath}: {e}")

    def process_file(self, filepath):
        """Processes the file based on its extension into .processed/"""
        filename = os.path.basename(filepath)

        # File size limit (H7): 50 MB
        try:
            if os.path.getsize(filepath) > 50 * 1024 * 1024:
                logger.warning(
                    f"[Watcher] File {filename} exceeds 50MB limit, skipping."
                )
                return
        except OSError:
            pass

        ext = os.path.splitext(filename)[1].lower()
        output_name = filename + ".txt"  # default extension
        output_path = os.path.join(self.processed_dir, output_name)

        from src.api.attachment_intake import is_vision_filename

        if is_vision_filename(filename):
            logger.info(
                "[Watcher] Skipping RAG processing for vision-only file %s", filename
            )
            return

        logger.info(f"[Watcher] Processing {filename} ({ext})...")

        try:
            # Document formats
            if ext == ".pdf":
                self._process_pdf(filepath, output_path)
            elif ext in [".csv", ".xlsx", ".xls", ".tsv"]:
                self._process_table(filepath, output_path, ext)
            elif ext in [".docx", ".doc"]:
                self._process_word(filepath, output_path, ext)
            elif ext == ".pptx":
                self._process_pptx(filepath, output_path)
            elif ext == ".epub":
                self._process_epub(filepath, output_path)
            elif ext == ".rtf":
                self._process_rtf(filepath, output_path)

            # Data/Config formats
            elif ext == ".json":
                self._process_json(filepath, output_path)
            elif ext == ".xml":
                self._process_xml(filepath, output_path)
            elif ext in [".yaml", ".yml"]:
                self._process_yaml(filepath, output_path)
            elif ext == ".toml":
                self._process_toml(filepath, output_path)
            elif ext in [".ini", ".conf", ".config"]:
                self._process_config(filepath, output_path)

            # Mark-up & Web
            elif ext == ".html" or ext == ".htm":
                self._process_html(filepath, output_path)
            elif ext == ".md" or ext == ".markdown":
                self._process_markdown(filepath, output_path)

            # Database & Archives
            elif ext == ".db" or ext == ".sqlite" or ext == ".sqlite3":
                self._process_sqlite(filepath, output_path)
            elif ext in [".zip", ".tar", ".gz", ".rar", ".7z"]:
                self._process_archive(filepath, output_path, ext)

            # Log files
            elif ext == ".log":
                self._process_log(filepath, output_path)

            # Source code (already readable, but can enhance)
            elif ext in [
                ".py",
                ".js",
                ".ts",
                ".java",
                ".cpp",
                ".c",
                ".go",
                ".rs",
                ".rb",
                ".php",
            ]:
                self._process_source_code(filepath, output_path, ext)

            else:
                # Try to read as plain text as fallback
                logger.info(f"[Watcher] Attempting to read {filename} as plain text...")
                self._process_plaintext(filepath, output_path)

            logger.info(f"[Watcher] Successfully processed {filename} -> {output_path}")
            if self.on_processed_callback:
                self.on_processed_callback(filepath, "processed")

        except Exception as e:
            logger.error(f"[Watcher] Failed to process {filename}: {e}")
            if self.on_processed_callback:
                self.on_processed_callback(filepath, "error")

    def _get_docling_converter(self):
        """Lazy-initialise Docling for DOCX only. Returns None if unavailable."""
        try:
            from docling.document_converter import (
                DocumentConverter,
                WordFormatOption,
            )
            from docling.datamodel.base_models import InputFormat
        except ImportError:
            logger.warning(
                "[Watcher] Docling not installed — falling back to python-docx for DOCX"
            )
            return None

        if self._docling_converter is None:
            self._docling_converter = DocumentConverter(
                format_options={
                    InputFormat.DOCX: WordFormatOption(),
                }
            )
            logger.info("[Watcher] Docling converter initialised (DOCX only)")
        return self._docling_converter

    def _process_pdf(self, filepath, output_path):
        """Extract PDF text via StirlingPDF (fallback PyMuPDF)."""
        from src.pdf.intake import extract_pdf_text_from_path

        text = extract_pdf_text_from_path(filepath, page_markers=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)

    def _process_table(self, filepath, output_path, ext):
        import pandas as pd

        def _clean_df(df: "pd.DataFrame") -> "pd.DataFrame":
            """Strip all-NaN rows/cols and fix Unnamed column headers."""
            df = (
                df.dropna(axis=0, how="all")
                .dropna(axis=1, how="all")
                .reset_index(drop=True)
            )
            if df.empty:
                return df

            # If column names are "Unnamed: N", infer from first data row
            if any(str(c).startswith("Unnamed:") for c in df.columns):
                first_row_non_empty = df.iloc[0].notna().any() if len(df) > 0 else False
                if first_row_non_empty and df.shape[0] > 1:
                    new_cols = []
                    for i, col in enumerate(df.columns):
                        cell_val = str(df.iloc[0, i]).strip() if i < df.shape[1] else ""
                        if cell_val and cell_val.lower() != "nan":
                            new_cols.append(cell_val)
                        else:
                            new_cols.append(df.columns[i])
                    df.columns = new_cols
                    df = df.iloc[1:].reset_index(drop=True)

                # Second pass: fix any remaining Unnamed: columns
                if any(str(c).startswith("Unnamed:") for c in df.columns):
                    for i, col in enumerate(df.columns):
                        if str(col).startswith("Unnamed:"):
                            found = False
                            for row_idx in range(min(5, len(df))):
                                cell_val = str(df.iloc[row_idx, i]).strip()
                                if (
                                    cell_val
                                    and cell_val.lower() != "nan"
                                    and not cell_val.startswith("Unnamed:")
                                ):
                                    df.rename(columns={col: cell_val}, inplace=True)
                                    found = True
                                    break
                            if not found:
                                df.rename(columns={col: f"Column_{i}"}, inplace=True)
                else:
                    df.columns = [f"Column_{i + 1}" for i in range(len(df.columns))]

            # Trim float noise: round floats to 4 significant digits
            for col in df.columns:
                if df[col].dtype == float:
                    df[col] = df[col].apply(lambda x: round(x, 4) if pd.notna(x) else x)
            return df

        output_path_md = output_path.replace(".txt", ".md")

        if ext == ".csv":
            df = pd.read_csv(filepath)
            df = _clean_df(df)
            with open(output_path_md, "w", encoding="utf-8") as f:
                f.write(df.to_markdown(index=False) if not df.empty else "(empty CSV)")
        else:
            # XLSX: read ALL sheets so context is not lost
            sheets: dict = pd.read_excel(filepath, sheet_name=None)
            parts: list[str] = []
            for sheet_name, df in sheets.items():
                df = _clean_df(df)
                if df.empty:
                    parts.append(f"## Sheet: {sheet_name}\n\n(empty sheet)")
                else:
                    rows, cols = df.shape
                    parts.append(
                        f"## Sheet: {sheet_name} ({rows} rows × {cols} columns)\n\n"
                        + df.to_markdown(index=False)
                    )
            combined = "\n\n---\n\n".join(parts) if parts else "(empty workbook)"
            with open(output_path_md, "w", encoding="utf-8") as f:
                f.write(combined)

    def _process_word(self, filepath, output_path, ext):
        """Extract DOC/DOCX text (using Docling/python-docx for DOCX, binary extractor for DOC)."""
        if ext == ".doc":
            from src.api.shared import extract_doc_text

            try:
                with open(filepath, "rb") as f:
                    raw_bytes = f.read()
                text = extract_doc_text(raw_bytes)
                if text:
                    with open(output_path, "w", encoding="utf-8") as f:
                        f.write(text)
                    return
            except Exception as e:
                logger.warning(
                    "[Watcher] Binary DOC extraction failed for %s: %s", filepath, e
                )
            return

        converter = self._get_docling_converter()
        if converter is not None:
            try:
                result = converter.convert(filepath)
                markdown_text = result.document.export_to_markdown()
                if markdown_text and len(markdown_text.strip()) > 50:
                    with open(output_path, "w", encoding="utf-8") as f:
                        f.write(markdown_text)
                    return
            except Exception as e:
                logger.warning(
                    "[Watcher] Docling DOCX extraction failed for %s: %s — falling back to python-docx",
                    filepath,
                    e,
                )

        # Fallback: python-docx (paragraphs only, no tables)
        from docx import Document

        doc = Document(filepath)
        text = ""
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"

        # Extract table content
        for table in doc.tables:
            text += "\n--- Table ---\n"
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                if row_idx == 0:
                    text += " | ".join(cells) + "\n" + "---" * len(cells) + "\n"
                else:
                    text += " | ".join(cells) + "\n"

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)

    def _process_json(self, filepath, output_path):
        """Extract and format JSON for readability."""
        import json

        try:
            with open(filepath, "r", encoding="utf-8") as f:
                data = json.load(f)

            # Pretty print with formatting
            formatted = json.dumps(data, indent=2, ensure_ascii=False)
            with open(output_path, "w", encoding="utf-8") as f:
                f.write("# JSON File\n\n```json\n" + formatted + "\n```")
        except json.JSONDecodeError as e:
            logger.warning(f"[Watcher] JSON decode error in {filepath}: {e}")
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# JSON Parse Error\n\nFailed to parse JSON: {e}\n\n")
                with open(filepath, "r", encoding="utf-8", errors="replace") as src:
                    f.write(src.read())

    def _process_xml(self, filepath, output_path):
        """Parse and nicely format XML."""
        try:
            import xml.etree.ElementTree as ET
            from xml.dom import minidom

            tree = ET.parse(filepath)
            xml_str = ET.tostring(tree.getroot(), encoding="unicode")

            # Pretty print XML
            dom = minidom.parseString(xml_str)
            pretty_xml = dom.toprettyxml(indent="  ")

            with open(output_path, "w", encoding="utf-8") as f:
                f.write("# XML File\n\n```xml\n" + pretty_xml + "\n```")
        except Exception as e:
            logger.warning(f"[Watcher] XML parse error in {filepath}: {e}")
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# XML Parse Error\n\nFailed to parse XML: {e}\n\n")
                with open(filepath, "r", encoding="utf-8", errors="replace") as src:
                    f.write(src.read())

    def _process_yaml(self, filepath, output_path):
        """Parse and format YAML."""
        try:
            import yaml

            with open(filepath, "r", encoding="utf-8") as f:
                data = yaml.safe_load(f)

            # Re-dump with nice formatting
            formatted = yaml.dump(data, default_flow_style=False, sort_keys=False)
            with open(output_path, "w", encoding="utf-8") as f:
                f.write("# YAML File\n\n```yaml\n" + formatted + "\n```")
        except ImportError:
            logger.warning("[Watcher] PyYAML not installed, treating as plain text")
            self._process_plaintext(filepath, output_path)
        except Exception as e:
            logger.warning(f"[Watcher] YAML parse error in {filepath}: {e}")
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# YAML Parse Error\n\nFailed to parse YAML: {e}\n\n")
                with open(filepath, "r", encoding="utf-8", errors="replace") as src:
                    f.write(src.read())

    def _process_toml(self, filepath, output_path):
        """Parse and format TOML."""
        try:
            try:
                import tomllib
            except ImportError:
                try:
                    import tomli as tomllib
                except ImportError:
                    raise ImportError("tomli or tomllib required for TOML parsing")

            with open(filepath, "rb") as f:
                data = tomllib.load(f)

            # Use TOML dumps if available, otherwise JSON format
            try:
                import tomli_w

                formatted = tomli_w.dumps(data)
            except ImportError:
                import json

                formatted = json.dumps(data, indent=2, ensure_ascii=False)

            with open(output_path, "w", encoding="utf-8") as f:
                f.write("# TOML File\n\n```toml\n" + formatted + "\n```")
        except ImportError:
            logger.warning(
                "[Watcher] TOML libraries not installed, treating as plain text"
            )
            self._process_plaintext(filepath, output_path)
        except Exception as e:
            logger.warning(f"[Watcher] TOML parse error in {filepath}: {e}")
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# TOML Parse Error\n\nFailed to parse TOML: {e}\n\n")
                with open(filepath, "r", encoding="utf-8", errors="replace") as src:
                    f.write(src.read())

    def _process_config(self, filepath, output_path):
        """Parse INI/CONF/CONFIG files."""
        try:
            import configparser

            config = configparser.ConfigParser()
            config.read(filepath, encoding="utf-8")

            with open(output_path, "w", encoding="utf-8") as f:
                f.write("# Configuration File\n\n")
                for section in config.sections():
                    f.write(f"## [{section}]\n")
                    for key, value in config.items(section):
                        f.write(f"- **{key}**: {value}\n")
                    f.write("\n")
        except Exception as e:
            logger.warning(f"[Watcher] Config parse error in {filepath}: {e}")
            self._process_plaintext(filepath, output_path)

    def _process_html(self, filepath, output_path):
        """Extract text and structure from HTML."""
        try:
            import re

            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                html_content = f.read()

            # Extract title if available
            title_match = re.search(
                r"<title[^>]*>([^<]+)</title>", html_content, re.IGNORECASE
            )
            title = title_match.group(1) if title_match else "HTML Document"

            # Extract text by removing script/style tags
            text = re.sub(
                r"<script[^>]*>.*?</script>",
                "",
                html_content,
                flags=re.IGNORECASE | re.DOTALL,
            )
            text = re.sub(
                r"<style[^>]*>.*?</style>", "", text, flags=re.IGNORECASE | re.DOTALL
            )

            # Remove HTML tags
            text = re.sub(r"<[^>]+>", "\n", text)
            # Clean up whitespace
            text = re.sub(r"\n\s*\n", "\n", text)
            text = text.strip()

            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# {title}\n\n")
                f.write("## Extracted Content\n\n")
                f.write(text)
        except Exception as e:
            logger.warning(f"[Watcher] HTML parsing error in {filepath}: {e}")
            self._process_plaintext(filepath, output_path)

    def _process_markdown(self, filepath, output_path):
        """Process Markdown - validate and ensure UTF-8."""
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                content = f.read()

            # Simply copy to .processed with validation
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(content)
        except UnicodeDecodeError:
            logger.warning(f"[Watcher] Markdown has encoding issues: {filepath}")
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(content)

    def _process_sqlite(self, filepath, output_path):
        """Extract schema and basic info from SQLite database."""
        try:
            import sqlite3

            conn = sqlite3.connect(filepath)
            cursor = conn.cursor()

            with open(output_path, "w", encoding="utf-8") as f:
                f.write("# SQLite Database\n\n")

                # Get all tables
                cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
                tables = cursor.fetchall()

                if not tables:
                    f.write("No tables found in database.\n")
                else:
                    f.write(f"## Tables ({len(tables)})\n\n")
                    for (table_name,) in tables:
                        f.write(f"### {table_name}\n")

                        # Get schema
                        cursor.execute(f'PRAGMA table_info("{table_name}")')
                        columns = cursor.fetchall()
                        f.write("**Columns:**\n")
                        for col in columns:
                            f.write(f"- {col[1]} ({col[2]})\n")

                        # Get row count
                        cursor.execute(f'SELECT COUNT(*) FROM "{table_name}"')
                        count = cursor.fetchone()[0]
                        f.write(f"\nRows: {count}\n\n")

                        # Sample data (first 5 rows)
                        if count > 0:
                            f.write("**Sample Data:**\n```\n")
                            cursor.execute(f'SELECT * FROM "{table_name}" LIMIT 5')
                            for row in cursor.fetchall():
                                f.write(str(row) + "\n")
                            f.write("```\n\n")

            conn.close()
        except Exception as e:
            logger.warning(f"[Watcher] SQLite parsing error in {filepath}: {e}")
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# SQLite Database Error\n\nFailed to parse database: {e}\n")

    def _process_archive(self, filepath, output_path, ext):
        """Extract archive contents listing."""
        try:
            if ext == ".zip":
                import zipfile

                with zipfile.ZipFile(filepath, "r") as archive:
                    files = archive.namelist()
            elif ext == ".tar" or ext == ".gz":
                import tarfile

                with tarfile.open(filepath, "r:*") as archive:
                    files = archive.getnames()
            elif ext in [".rar", ".7z"]:
                logger.warning(
                    f"[Watcher] {ext} requires additional libraries, skipping detailed processing"
                )
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(
                        f"# {ext} Archive\n\nSupported but requires additional libraries for detailed extraction.\n"
                    )
                return
            else:
                files = []

            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# Archive Contents ({ext})\n\n")
                f.write(f"**Total items:** {len(files)}\n\n")
                f.write("## File List\n\n")
                for fname in sorted(files)[
                    :100
                ]:  # Limit to first 100 to avoid huge outputs
                    f.write(f"- {fname}\n")
                if len(files) > 100:
                    f.write(f"\n... and {len(files) - 100} more items\n")
        except Exception as e:
            logger.warning(f"[Watcher] Archive parsing error in {filepath}: {e}")
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# Archive Error\n\nFailed to process archive: {e}\n")

    def _process_log(self, filepath, output_path):
        """Process log files - tail and summarize."""
        try:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                lines = f.readlines()

            total_lines = len(lines)
            # Get last 500 lines or all if less
            tail_lines = lines[-500:] if total_lines > 500 else lines

            with open(output_path, "w", encoding="utf-8") as f:
                f.write("# Log File\n\n")
                f.write(f"**Total lines:** {total_lines}\n")
                f.write(f"**Showing:** Last {len(tail_lines)} lines\n\n")
                f.write("## Content\n\n```\n")
                f.writelines(tail_lines)
                f.write("```\n")
        except Exception as e:
            logger.warning(f"[Watcher] Log parsing error in {filepath}: {e}")
            self._process_plaintext(filepath, output_path)

    def _process_source_code(self, filepath, output_path, ext):
        """Process source code files with metadata."""
        try:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()

            # Count lines and estimate functions/classes
            lines = content.split("\n")
            line_count = len(lines)

            # Simple heuristic for functions/classes
            import re

            if ext == ".py":
                func_pattern = r"^\s*def\s+\w+"
                class_pattern = r"^\s*class\s+\w+"
            elif ext in [".js", ".ts"]:
                func_pattern = r"^\s*(async\s+)?function\s+\w+|const\s+\w+\s*=.*=>"
                class_pattern = r"^\s*class\s+\w+"
            elif ext in [".java"]:
                func_pattern = r"^\s*public\s+\w+\s+\w+\s*\("
                class_pattern = r"^\s*public\s+class\s+\w+"
            else:
                func_pattern = None
                class_pattern = None

            func_count = (
                len(re.findall(func_pattern, content, re.MULTILINE))
                if func_pattern
                else 0
            )
            class_count = (
                len(re.findall(class_pattern, content, re.MULTILINE))
                if class_pattern
                else 0
            )

            # Extract first docstring/comment if present
            first_doc = ""
            if ext == ".py":
                doc_match = re.search(r'"""(.*?)"""', content, re.DOTALL)
                if doc_match:
                    first_doc = doc_match.group(1).strip()[:200]

            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# Source Code ({ext[1:].upper()})\n\n")
                f.write(f"**Lines:** {line_count}\n")
                f.write(f"**Functions/Methods:** {func_count}\n")
                f.write(f"**Classes:** {class_count}\n")
                if first_doc:
                    f.write(f"**Description:** {first_doc}\n")
                f.write(f"\n## Code\n\n```{ext[1:]}\n")
                f.write(content)
                f.write("\n```")
        except Exception as e:
            logger.warning(f"[Watcher] Source code parsing error in {filepath}: {e}")
            self._process_plaintext(filepath, output_path)

    def _process_pptx(self, filepath, output_path):
        """Extract text from PowerPoint presentations."""
        try:
            from pptx import Presentation

            prs = Presentation(filepath)
            lines = [f"# PowerPoint: {os.path.basename(filepath)}\n"]
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"\n## Slide {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                lines.append(text)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            lines.append(" | ".join(cells))
            with open(output_path, "w", encoding="utf-8") as f:
                f.write("\n".join(lines))
        except ImportError:
            logger.warning("[Watcher] python-pptx not installed, skipping pptx")
            self._process_plaintext(filepath, output_path)

    def _process_epub(self, filepath, output_path):
        """Extract text from EPUB ebooks."""
        try:
            import zipfile
            from bs4 import BeautifulSoup

            lines = [f"# EPUB: {os.path.basename(filepath)}\n"]
            with zipfile.ZipFile(filepath, "r") as zf:
                for name in zf.namelist():
                    if name.endswith((".xhtml", ".html", ".htm")):
                        try:
                            html = zf.read(name).decode("utf-8", errors="replace")
                            soup = BeautifulSoup(html, "html.parser")
                            for tag in soup(["script", "style", "nav"]):
                                tag.decompose()
                            text = soup.get_text(separator="\n", strip=True)
                            if text.strip():
                                lines.append(f"\n## {name}\n")
                                lines.append(text)
                        except Exception as e:
                            logger.warning("Error suppressed: %s", e)
                            continue
            with open(output_path, "w", encoding="utf-8") as f:
                f.write("\n".join(lines))
        except Exception as e:
            logger.warning(f"[Watcher] EPUB processing error: {e}")
            self._process_plaintext(filepath, output_path)

    def _process_rtf(self, filepath, output_path):
        """Extract text from RTF files."""
        try:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            import re

            text = re.sub(r"\\[a-z]+\d*\s?", "", content)
            text = re.sub(r"[{}]", "", text)
            text = text.strip()
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# RTF: {os.path.basename(filepath)}\n\n{text}")
        except Exception as e:
            logger.warning(f"[Watcher] RTF processing error: {e}")
            self._process_plaintext(filepath, output_path)

    def _process_plaintext(self, filepath, output_path):
        """Fallback: read file as plain text."""
        try:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()

            # Detect if it might be a code file even without extension
            if len(content) < 100000:  # Only for reasonably sized files
                import re

                if re.search(
                    r"(function|class|def |import |include|struct|interface)",
                    content,
                    re.IGNORECASE,
                ):
                    # Likely code
                    with open(output_path, "w", encoding="utf-8") as f:
                        f.write("# Unrecognized Code File\n\n```\n" + content + "\n```")
                else:
                    with open(output_path, "w", encoding="utf-8") as f:
                        f.write(content)
            else:
                # Very large file - just write as-is
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(content)
        except Exception as e:
            logger.error(f"[Watcher] Failed to read {filepath} as plain text: {e}")
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# Error Reading File\n\nCould not read file: {e}\n")


def _catchup_unprocessed_files(workspace_dir, handler):
    """Scan workspace for files modified while the watcher was offline."""
    logger.info(f"[Watcher] Scanning for offline modifications in {workspace_dir}...")
    try:
        count = 0
        for root, dirs, files in os.walk(workspace_dir):
            # Skip hidden directories and .processed cache
            dirs[:] = [d for d in dirs if not d.startswith(".") and d != ".processed"]

            for filename in files:
                if filename.startswith("."):
                    continue

                filepath = os.path.join(root, filename)
                try:
                    src_mtime = os.path.getmtime(filepath)

                    # Determine expected processed file path
                    # By default it's filename + .txt
                    expected_txt = os.path.join(
                        handler.processed_dir, filename + ".txt"
                    )

                    # For tables, it's filename + .txt replaced with .md in the code,
                    # which effectively means filename + .md for files that originally end with .csv, etc?
                    # Wait, the code does: output_path.replace('.txt', '.md')
                    # So if output_path is `filename + '.txt'`, then it becomes `filename + '.md'`
                    expected_md = os.path.join(handler.processed_dir, filename + ".md")

                    needs_processing = False

                    # Check if either txt or md exists
                    if os.path.exists(expected_txt):
                        proc_mtime = os.path.getmtime(expected_txt)
                        if src_mtime > proc_mtime:
                            needs_processing = True
                    elif os.path.exists(expected_md):
                        proc_mtime = os.path.getmtime(expected_md)
                        if src_mtime > proc_mtime:
                            needs_processing = True
                    else:
                        # Neither exists
                        needs_processing = True

                    if needs_processing:
                        logger.info(
                            f"[Watcher] Catch-up processing required for: {filepath}"
                        )
                        # Process synchronously to avoid spawning thousands of threads on first boot
                        handler.process_file(filepath)
                        count += 1
                except Exception as e:
                    logger.warning(
                        f"[Watcher] Error checking mtime for {filepath}: {e}"
                    )

        if count > 0:
            logger.info(f"[Watcher] Catch-up complete. Processed {count} missed files.")
        else:
            logger.info("[Watcher] Catch-up complete. No missed files found.")
    except Exception as e:
        logger.error(f"[Watcher] Error during offline catch-up scan: {e}")


def start_watcher(workspace_dir, on_processed_callback=None):
    """Start the watchdog observer thread in the background.

    Args:
        workspace_dir: Path to the workspace directory to monitor.
        on_processed_callback: Optional callback ``(filename, status)``
            invoked after each file is processed.

    Returns:
        The running ``Observer`` instance (call ``.stop()`` to shut down).
    """
    workspace_dir = os.path.abspath(workspace_dir)
    event_handler = FileWatcherHandler(
        workspace_dir, on_processed_callback=on_processed_callback
    )

    # Sweep for files that were added/modified while the server was offline
    _catchup_unprocessed_files(workspace_dir, event_handler)

    observer = Observer()
    observer.schedule(event_handler, path=workspace_dir, recursive=True)
    observer.start()
    logger.info(f"🚀 Started WorkspaceWatcher on {workspace_dir}")
    return observer


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument("--workspace", default="./workspace")
    args = parser.parse_args()

    observer = start_watcher(args.workspace)
    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
    observer.join()
